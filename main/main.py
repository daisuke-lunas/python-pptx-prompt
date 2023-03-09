# encoding: utf-8

import collections.abc
# ↑↑↑のimportの理由は、https://kakakakakku.hatenablog.com/entry/2022/11/24/105451 参照。
# collectionsクラスをPythonが非推奨にした影響の解決法

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.connector import Connector
import xml.etree.ElementTree as ET
from connectorDetail import ConnectorDetail
import lxml.etree as etree

prs = Presentation("test.2.pptx")
shapeList: list[Shape] = []
connectorList: list[Connector] = []

for slide in prs.slides:
    for shape in slide.shapes:
        if isinstance(shape, Shape):
            shapeList.append(shape)
            print("SHAPE: id&text: " + str(shape.shape_id) + " " + shape.text)
            print("top: " + str(shape.top))
            print("left: " + str(shape.left))
            print("\n")
            continue

        elif isinstance(shape, Connector):
            connectorList.append(shape)
            print("CONNECTOR: shape_id: " + str(shape.shape_id))
            print("begin: (" + str(shape.begin_x) +
                  "," + str(shape.begin_y) + ")")
            print("end: (" + str(shape.end_x) +
                  "," + str(shape.end_y) + ")")
            print("height: " + str(shape.height))
            print("width: " + str(shape.width))
            print("\n")
            continue
        else:
            print("other shape detected!")
            print(shape)
            print("\n")

connectorDtlList: list[ConnectorDetail] = []

for connector in connectorList:
    # connectorの先端の形状を判定し、start/endを決める
    conn_elm: etree.ElementBase = connector.element
    for desc in conn_elm.iterdescendants():
        desc: etree.ElementBase = desc

        if desc.get("type") == "triangle":
            print("tag name is: " + desc.tag)
            # 今回はtailEndのみ、末端が矢印頭
            dtl = ConnectorDetail()
            dtl.connector = connector
            # shapeListをイテレート
            for shp in shapeList:
                if dtl.connectorStartsAt(shp):
                    dtl.startAt = shp
                if dtl.connectorEndAt(shp):
                    dtl.endAt = shp
                # begin_x/yを見て、shapeの中/外側周辺にあるかを判定
                # あったら、startAtに入れる
                # end_x/yも同じく
            connectorDtlList.append(dtl)

for dtl in connectorDtlList:
    # 期待結果のCSVを作成する
    # [startAtのtext],[endAtのtext]でコンソールに
    id = str(dtl.connector.shape_id)
    startText = dtl.startAt.text if dtl.startAt != None else "なし"
    endText = dtl.endAt.text if dtl.endAt != None else "なし"
    print(id, ",", startText, ",", endText)

print("finish!")
